from selenium import webdriver
import time
from PIL import Image
from io import BytesIO
import cv2
import numpy as np
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from pptx import Presentation 
from pptx.util import Inches
 
X_OFFSET = 379
Y_OFFSET = -123
 
def replace_image_in_cropped_area(screenshot, reference_image_path, img_location):
    """Replace the image in the specified location of the screenshot with the reference image."""
    reference_image = Image.open(reference_image_path)
 
    # Get the dimensions of the reference image
    reference_width, reference_height = reference_image.size
 
    # Resize the reference image to match the original image size
    reference_image = reference_image.resize((img_location[2], img_location[3]))
 
    # Calculate the precise pasting location with an offset
    paste_location = (img_location[0] + X_OFFSET, img_location[1] + Y_OFFSET)
 
    # Paste the reference image into the specified location
    screenshot.paste(reference_image, paste_location)
#     time.sleep(3)
 
    # Save the modified screenshot
    modified_screenshot_path = r"C:\Ad_screenshot\modified_screenshot.png"
    screenshot.save(modified_screenshot_path)
 
    # Display the modified screenshot (optional)
    screenshot.show()
 
    return modified_screenshot_path
 
def find_and_replace_reference_image(driver, reference_image_path):
    # Load the reference image
    reference_image = cv2.imread(reference_image_path, cv2.IMREAD_UNCHANGED)
    reference_height, reference_width, _ = reference_image.shape
 
    images = driver.find_elements(By.TAG_NAME, 'iframe')
 
    # Capture a screenshot for the entire webpage
    screenshot = driver.get_screenshot_as_png()
    screenshot = Image.open(BytesIO(screenshot))
    screenshot = np.array(screenshot)
 
    for img in images:
        # Get the dimensions of the current image
        img_width = img.size['width']
        img_height = img.size['height']
 
        # Print the dimensions of the current image for debugging
        print(f"Image Dimensions: {img_width} x {img_height}")
 
        # Compare dimensions with the reference image
        if img_width == reference_width and img_height == reference_height:
            # Scroll to the element
            driver.execute_script("arguments[0].scrollIntoView({behavior: 'auto', block: 'center'});", img)
            time.sleep(0)  # Adjust the sleep time if needed
 
            # Capture a screenshot for the entire webpage
            screenshot = driver.get_screenshot_as_png()
            screenshot = Image.open(BytesIO(screenshot))
            screenshot = np.array(screenshot)
 
            # Convert images to the appropriate data type
            screenshot_gray = cv2.cvtColor(screenshot, cv2.COLOR_BGR2GRAY)
            reference_gray = cv2.cvtColor(reference_image, cv2.COLOR_BGR2GRAY)
 
            # Perform template matching
            result = cv2.matchTemplate(screenshot_gray, reference_gray, cv2.TM_CCOEFF_NORMED)
            _, _, _, max_loc = cv2.minMaxLoc(result)
 
            # Replace the reference image at the location of the matched template
            modified_screenshot_path = replace_image_in_cropped_area(
                Image.fromarray(screenshot),
                reference_image_path,
                (*max_loc, reference_width, reference_height)
            )
 
            return modified_screenshot_path
 
    print("No image found with the dimensions of the reference image.")
    return None
 
def save_screenshot_to_pptx(screenshot_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Use the layout that suits your needs
 
    slide = prs.slides.add_slide(slide_layout)
 
    # Add the screenshot as a static image
    left = Inches(-0.8)
    top = Inches(0)
    img = Image.open(screenshot_path)
    width, height = img.size
    img_path =  r"C:\Ad_screenshot\static_screenshot.png"
    img.save(img_path)
    slide.shapes.add_picture(img_path, left, top, width=Inches(width / 80), height=Inches(height / 80))
 
    pptx_path = r"C:\Ad_screenshot\modified_screenshot_300_250.pptx"
    prs.save(pptx_path)
    
if __name__ == "__main__":
    url = "https://www.floridatravellife.com/florida-outdoor-activities/best-natural-springs-in-florida/"
    reference_image_path =  r"C:\Ad_screenshot\Image_sojern\downloaded_image_sojern_300_250.jpg"
 
    driver = webdriver.Chrome()
    driver.get(url)
    driver.maximize_window()
    time.sleep(20)
 
    # Find and replace the reference image at the location of the first image with the same dimensions
    modified_screenshot_path = find_and_replace_reference_image(driver, reference_image_path)
 
    if modified_screenshot_path:
        print(f"Modified screenshot saved at: {modified_screenshot_path}")
        # Save the modified screenshot to PowerPoint
        save_screenshot_to_pptx(modified_screenshot_path)
    # Close the browser
    driver.quit()